"""
Enhanced document parser supporting multimodal content extraction.
Inspired by RAG-Anything's MineruParser / DoclingParser architecture.

This parser extracts text, images, tables, and equations from documents,
producing a unified content_list format.
"""

import os
import re
import json
import traceback
from pathlib import Path
from typing import Dict, List, Optional, Tuple

try:
    import fitz  # PyMuPDF
    PYMUPDF_AVAILABLE = True
except ImportError:
    PYMUPDF_AVAILABLE = False

try:
    from docx import Document as DocxDocument
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    from PIL import Image
    PIL_AVAILABLE = True
except ImportError:
    PIL_AVAILABLE = False


class DocumentParserService:
    """
    Enhanced document parser that extracts structured content including
    text, images, tables, and equations from various document formats.

    Produces a standardized content_list format:
    [
        {"type": "text", "text": "...", "page_idx": 0},
        {"type": "image", "img_path": "/path/to/img.png", "image_caption": [...], "page_idx": 1},
        {"type": "table", "table_body": "markdown", "table_caption": [...], "page_idx": 2},
        {"type": "equation", "text": "LaTeX", "equation_format": "latex", "page_idx": 3},
    ]
    """

    # Supported format groups
    PDF_FORMATS = {".pdf"}
    IMAGE_FORMATS = {".jpg", ".jpeg", ".png", ".bmp", ".tiff", ".tif", ".gif", ".webp"}
    OFFICE_FORMATS = {".docx", ".doc", ".pptx", ".ppt", ".xlsx", ".xls"}
    TEXT_FORMATS = {".txt", ".md", ".csv", ".html", ".htm"}

    def __init__(self, output_dir: str = "./data/parser_output"):
        self.output_dir = Path(output_dir)
        self.output_dir.mkdir(parents=True, exist_ok=True)

    def parse_document(self, file_path: str) -> Dict:
        """
        Parse a document and return structured content.

        Returns:
            {
                "content_list": [...],  # Standardized content blocks
                "text": str,           # Extracted full text
                "chunks": [...],       # Text chunks for RAG
                "metadata": {...},     # File metadata
                "multimodal_count": {  # Count of each content type
                    "text": N, "image": N, "table": N, "equation": N
                }
            }
        """
        path = Path(file_path)
        ext = path.suffix.lower()

        if ext in self.PDF_FORMATS:
            content_list = self._parse_pdf(file_path)
        elif ext in self.IMAGE_FORMATS:
            content_list = self._parse_image(file_path)
        elif ext in self.OFFICE_FORMATS:
            content_list = self._parse_office(file_path)
        elif ext in self.TEXT_FORMATS:
            content_list = self._parse_text(file_path)
        else:
            raise ValueError(f"Unsupported format: {ext}")

        # Extract text and compute chunks
        from .utils import separate_content, chunk_text
        text_content, multimodal_items = separate_content(content_list)
        chunks = chunk_text(text_content)

        # Count content types
        type_counts = {"text": 0, "image": 0, "table": 0, "equation": 0}
        for item in content_list:
            t = item.get("type", "text")
            if t in type_counts:
                type_counts[t] += 1
            else:
                type_counts[t] = type_counts.get(t, 0) + 1

        return {
            "content_list": content_list,
            "text": text_content,
            "chunks": chunks,
            "metadata": {
                "filename": path.name,
                "file_type": ext,
                "file_size_mb": round(path.stat().st_size / (1024 * 1024), 2),
            },
            "multimodal_count": type_counts,
            "num_chunks": len(chunks),
            "multimodal_items": multimodal_items,
        }

    def _parse_pdf(self, file_path: str) -> List[Dict]:
        """Parse PDF with text, image, and table extraction."""
        if not PYMUPDF_AVAILABLE:
            return [{"type": "text", "text": "PyMuPDF not available", "page_idx": 0}]

        content_list = []
        images_dir = self.output_dir / Path(file_path).stem / "images"
        images_dir.mkdir(parents=True, exist_ok=True)

        try:
            doc = fitz.open(file_path)

            for page_idx, page in enumerate(doc):
                # Extract text blocks with position info
                blocks = page.get_text("dict", flags=fitz.TEXT_PRESERVE_WHITESPACE)["blocks"]

                page_text_parts = []

                for block_idx, block in enumerate(blocks):
                    if block["type"] == 0:  # Text block
                        text = ""
                        for line in block.get("lines", []):
                            for span in line.get("spans", []):
                                text += span.get("text", "")
                            text += "\n"
                        text = text.strip()

                        if text:
                            # Check if it looks like a table (has tab-separated data)
                            if self._looks_like_table(text):
                                content_list.append({
                                    "type": "table",
                                    "table_body": self._text_to_markdown_table(text),
                                    "table_caption": [],
                                    "table_footnote": [],
                                    "page_idx": page_idx,
                                })
                            # Check if it looks like an equation
                            elif self._looks_like_equation(text):
                                content_list.append({
                                    "type": "equation",
                                    "text": text,
                                    "equation_format": "text",
                                    "page_idx": page_idx,
                                })
                            else:
                                page_text_parts.append(text)

                    elif block["type"] == 1:  # Image block
                        try:
                            img_data = block.get("image", b"")
                            if img_data and len(img_data) > 100:  # Skip tiny images
                                img_filename = f"page{page_idx}_img{block_idx}.png"
                                img_path = str(images_dir / img_filename)

                                with open(img_path, "wb") as f:
                                    f.write(img_data)

                                content_list.append({
                                    "type": "image",
                                    "img_path": img_path,
                                    "image_caption": [],
                                    "image_footnote": [],
                                    "page_idx": page_idx,
                                })
                        except Exception as e:
                            print(f"Error extracting image from page {page_idx}: {e}")

                # Add accumulated text for this page
                if page_text_parts:
                    combined_text = "\n".join(page_text_parts)
                    if combined_text.strip():
                        content_list.append({
                            "type": "text",
                            "text": combined_text,
                            "page_idx": page_idx,
                        })

            # Also extract embedded images via xref
            for page_idx, page in enumerate(doc):
                image_list = page.get_images(full=True)
                for img_idx, img_info in enumerate(image_list):
                    xref = img_info[0]
                    try:
                        pix = fitz.Pixmap(doc, xref)
                        if pix.n < 5:  # GRAY or RGB
                            img_filename = f"page{page_idx}_xref{xref}.png"
                        else:  # CMYK or other
                            pix = fitz.Pixmap(fitz.csRGB, pix)
                            img_filename = f"page{page_idx}_xref{xref}.png"

                        img_path = str(images_dir / img_filename)
                        # Only save if not too small
                        if pix.width > 50 and pix.height > 50:
                            pix.save(img_path)
                            # Check if this image was already extracted via blocks
                            already_exists = any(
                                item.get("img_path", "").endswith(img_filename)
                                for item in content_list
                                if item.get("type") == "image"
                            )
                            if not already_exists:
                                content_list.append({
                                    "type": "image",
                                    "img_path": img_path,
                                    "image_caption": [],
                                    "image_footnote": [],
                                    "page_idx": page_idx,
                                })
                        pix = None
                    except Exception:
                        pass

            doc.close()

        except Exception as e:
            print(f"Error parsing PDF {file_path}: {e}")
            traceback.print_exc()
            # Fallback: simple text extraction
            try:
                doc = fitz.open(file_path)
                text = ""
                for page in doc:
                    text += page.get_text()
                doc.close()
                if text.strip():
                    content_list = [{"type": "text", "text": text, "page_idx": 0}]
            except Exception:
                content_list = [{"type": "text", "text": f"Error parsing PDF: {e}", "page_idx": 0}]

        return content_list

    def _parse_image(self, file_path: str) -> List[Dict]:
        """Parse a standalone image file."""
        return [{
            "type": "image",
            "img_path": str(Path(file_path).resolve()),
            "image_caption": [],
            "image_footnote": [],
            "page_idx": 0,
        }]

    def _parse_office(self, file_path: str) -> List[Dict]:
        """Parse Office documents (DOCX, XLSX, PPTX)."""
        ext = Path(file_path).suffix.lower()
        content_list = []

        if ext in (".docx", ".doc"):
            content_list = self._parse_docx(file_path)
        elif ext in (".xlsx", ".xls"):
            content_list = self._parse_xlsx(file_path)
        elif ext in (".pptx", ".ppt"):
            content_list = self._parse_pptx(file_path)
        else:
            content_list = [{"type": "text", "text": f"Unsupported office format: {ext}", "page_idx": 0}]

        return content_list

    def _parse_docx(self, file_path: str) -> List[Dict]:
        """Parse DOCX with text, images, and tables."""
        if not DOCX_AVAILABLE:
            return [{"type": "text", "text": "python-docx not available", "page_idx": 0}]

        content_list = []
        images_dir = self.output_dir / Path(file_path).stem / "images"
        images_dir.mkdir(parents=True, exist_ok=True)

        try:
            doc = DocxDocument(file_path)

            # Extract paragraphs
            for i, para in enumerate(doc.paragraphs):
                text = para.text.strip()
                if text:
                    content_list.append({
                        "type": "text",
                        "text": text,
                        "page_idx": 0,
                    })

            # Extract tables
            for i, table in enumerate(doc.tables):
                rows = []
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    rows.append(cells)

                if rows:
                    # Convert to markdown table
                    md_table = self._rows_to_markdown_table(rows)
                    content_list.append({
                        "type": "table",
                        "table_body": md_table,
                        "table_caption": [],
                        "table_footnote": [],
                        "page_idx": 0,
                    })

            # Extract images from relationships
            img_count = 0
            for rel in doc.part.rels.values():
                if "image" in rel.reltype:
                    try:
                        img_data = rel.target_part.blob
                        img_ext = Path(rel.target_ref).suffix or ".png"
                        img_filename = f"docx_img_{img_count}{img_ext}"
                        img_path = str(images_dir / img_filename)
                        with open(img_path, "wb") as f:
                            f.write(img_data)
                        content_list.append({
                            "type": "image",
                            "img_path": img_path,
                            "image_caption": [],
                            "image_footnote": [],
                            "page_idx": 0,
                        })
                        img_count += 1
                    except Exception as e:
                        print(f"Error extracting DOCX image: {e}")

        except Exception as e:
            print(f"Error parsing DOCX {file_path}: {e}")
            content_list = [{"type": "text", "text": f"Error parsing DOCX: {e}", "page_idx": 0}]

        return content_list

    def _parse_xlsx(self, file_path: str) -> List[Dict]:
        """Parse Excel files as tables."""
        content_list = []
        try:
            import openpyxl
            wb = openpyxl.load_workbook(file_path, data_only=True)

            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                rows = []
                for row in ws.iter_rows(values_only=True):
                    cells = [str(cell) if cell is not None else "" for cell in row]
                    if any(c.strip() for c in cells):
                        rows.append(cells)

                if rows:
                    md_table = self._rows_to_markdown_table(rows)
                    content_list.append({
                        "type": "table",
                        "table_body": md_table,
                        "table_caption": [f"Sheet: {sheet_name}"],
                        "table_footnote": [],
                        "page_idx": 0,
                    })

        except ImportError:
            content_list = [{"type": "text", "text": "openpyxl not installed for Excel parsing", "page_idx": 0}]
        except Exception as e:
            content_list = [{"type": "text", "text": f"Error parsing Excel: {e}", "page_idx": 0}]

        return content_list

    def _parse_pptx(self, file_path: str) -> List[Dict]:
        """Parse PowerPoint files."""
        content_list = []
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            images_dir = self.output_dir / Path(file_path).stem / "images"
            images_dir.mkdir(parents=True, exist_ok=True)

            for slide_idx, slide in enumerate(prs.slides):
                slide_texts = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            text = para.text.strip()
                            if text:
                                slide_texts.append(text)

                    if shape.has_table:
                        rows = []
                        for row in shape.table.rows:
                            cells = [cell.text.strip() for cell in row.cells]
                            rows.append(cells)
                        if rows:
                            md_table = self._rows_to_markdown_table(rows)
                            content_list.append({
                                "type": "table",
                                "table_body": md_table,
                                "table_caption": [f"Slide {slide_idx + 1}"],
                                "table_footnote": [],
                                "page_idx": slide_idx,
                            })

                    if shape.shape_type == 13:  # Picture
                        try:
                            img_data = shape.image.blob
                            img_ext = shape.image.content_type.split("/")[-1]
                            img_filename = f"slide{slide_idx}_img.{img_ext}"
                            img_path = str(images_dir / img_filename)
                            with open(img_path, "wb") as f:
                                f.write(img_data)
                            content_list.append({
                                "type": "image",
                                "img_path": img_path,
                                "image_caption": [],
                                "image_footnote": [],
                                "page_idx": slide_idx,
                            })
                        except Exception:
                            pass

                if slide_texts:
                    content_list.append({
                        "type": "text",
                        "text": "\n".join(slide_texts),
                        "page_idx": slide_idx,
                    })

        except ImportError:
            content_list = [{"type": "text", "text": "python-pptx not installed", "page_idx": 0}]
        except Exception as e:
            content_list = [{"type": "text", "text": f"Error parsing PPTX: {e}", "page_idx": 0}]

        return content_list

    def _parse_text(self, file_path: str) -> List[Dict]:
        """Parse plain text, markdown, HTML, or CSV files."""
        ext = Path(file_path).suffix.lower()
        content_list = []

        try:
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                text = f.read()

            if ext == ".csv":
                content_list = self._parse_csv_content(text, file_path)
            elif ext in (".html", ".htm"):
                content_list = self._parse_html_content(text)
            elif ext == ".md":
                content_list = self._parse_markdown_content(text)
            else:
                content_list = [{"type": "text", "text": text, "page_idx": 0}]

        except Exception as e:
            content_list = [{"type": "text", "text": f"Error reading file: {e}", "page_idx": 0}]

        return content_list

    def _parse_csv_content(self, text: str, file_path: str) -> List[Dict]:
        """Parse CSV content as a table."""
        try:
            import csv
            import io
            reader = csv.reader(io.StringIO(text))
            rows = [row for row in reader if any(cell.strip() for cell in row)]
            if rows:
                md_table = self._rows_to_markdown_table(rows)
                return [{
                    "type": "table",
                    "table_body": md_table,
                    "table_caption": [Path(file_path).name],
                    "table_footnote": [],
                    "page_idx": 0,
                }]
        except Exception:
            pass
        return [{"type": "text", "text": text, "page_idx": 0}]

    def _parse_html_content(self, text: str) -> List[Dict]:
        """Parse HTML content, extract text and tables."""
        content_list = []
        try:
            # Simple regex-based extraction (avoid heavy BeautifulSoup dependency)
            # Remove script and style
            text = re.sub(r"<script[^>]*>.*?</script>", "", text, flags=re.DOTALL)
            text = re.sub(r"<style[^>]*>.*?</style>", "", text, flags=re.DOTALL)

            # Extract tables
            tables = re.findall(r"<table[^>]*>(.*?)</table>", text, re.DOTALL)
            for table_html in tables:
                rows = []
                tr_matches = re.findall(r"<tr[^>]*>(.*?)</tr>", table_html, re.DOTALL)
                for tr in tr_matches:
                    cells = re.findall(r"<t[dh][^>]*>(.*?)</t[dh]>", tr, re.DOTALL)
                    cells = [re.sub(r"<[^>]+>", "", c).strip() for c in cells]
                    if cells:
                        rows.append(cells)
                if rows:
                    md_table = self._rows_to_markdown_table(rows)
                    content_list.append({
                        "type": "table",
                        "table_body": md_table,
                        "table_caption": [],
                        "table_footnote": [],
                        "page_idx": 0,
                    })

            # Remove tags and get plain text
            plain_text = re.sub(r"<[^>]+>", " ", text)
            plain_text = re.sub(r"\s+", " ", plain_text).strip()
            if plain_text:
                content_list.append({"type": "text", "text": plain_text, "page_idx": 0})

        except Exception:
            content_list = [{"type": "text", "text": text, "page_idx": 0}]

        return content_list

    def _parse_markdown_content(self, text: str) -> List[Dict]:
        """Parse Markdown content, detect tables and equations."""
        content_list = []
        lines = text.split("\n")
        current_text = []
        i = 0

        while i < len(lines):
            line = lines[i]

            # Detect markdown table
            if "|" in line and i + 1 < len(lines) and re.match(r"\s*\|[\s\-:|]+\|\s*$", lines[i + 1]):
                # Flush current text
                if current_text:
                    content_list.append({
                        "type": "text",
                        "text": "\n".join(current_text).strip(),
                        "page_idx": 0,
                    })
                    current_text = []

                # Collect table lines
                table_lines = [line]
                i += 1
                while i < len(lines) and "|" in lines[i]:
                    table_lines.append(lines[i])
                    i += 1

                content_list.append({
                    "type": "table",
                    "table_body": "\n".join(table_lines),
                    "table_caption": [],
                    "table_footnote": [],
                    "page_idx": 0,
                })
                continue

            # Detect LaTeX equations ($$...$$)
            if line.strip().startswith("$$"):
                if current_text:
                    content_list.append({
                        "type": "text",
                        "text": "\n".join(current_text).strip(),
                        "page_idx": 0,
                    })
                    current_text = []

                eq_lines = [line]
                if not line.strip().endswith("$$") or line.strip() == "$$":
                    i += 1
                    while i < len(lines) and "$$" not in lines[i]:
                        eq_lines.append(lines[i])
                        i += 1
                    if i < len(lines):
                        eq_lines.append(lines[i])

                eq_text = "\n".join(eq_lines).replace("$$", "").strip()
                if eq_text:
                    content_list.append({
                        "type": "equation",
                        "text": eq_text,
                        "equation_format": "latex",
                        "page_idx": 0,
                    })
                i += 1
                continue

            # Detect image references ![alt](path)
            img_match = re.match(r"!\[([^\]]*)\]\(([^)]+)\)", line.strip())
            if img_match:
                if current_text:
                    content_list.append({
                        "type": "text",
                        "text": "\n".join(current_text).strip(),
                        "page_idx": 0,
                    })
                    current_text = []

                alt_text = img_match.group(1)
                img_path = img_match.group(2)
                content_list.append({
                    "type": "image",
                    "img_path": img_path,
                    "image_caption": [alt_text] if alt_text else [],
                    "image_footnote": [],
                    "page_idx": 0,
                })
                i += 1
                continue

            current_text.append(line)
            i += 1

        # Flush remaining text
        if current_text:
            final_text = "\n".join(current_text).strip()
            if final_text:
                content_list.append({
                    "type": "text",
                    "text": final_text,
                    "page_idx": 0,
                })

        return content_list

    # ---- Utility Methods ----

    def _looks_like_table(self, text: str) -> bool:
        """Check if text block looks like tabular data."""
        lines = text.strip().split("\n")
        if len(lines) < 2:
            return False
        # Check for consistent tab/pipe delimiters
        tab_counts = [line.count("\t") for line in lines]
        if tab_counts[0] > 1 and all(c == tab_counts[0] for c in tab_counts[:3]):
            return True
        pipe_counts = [line.count("|") for line in lines]
        if pipe_counts[0] > 1 and all(c == pipe_counts[0] for c in pipe_counts[:3]):
            return True
        return False

    def _looks_like_equation(self, text: str) -> bool:
        """Check if text block looks like a mathematical equation."""
        text = text.strip()
        if not text:
            return False
        # LaTeX markers
        if text.startswith("\\") or "\\frac" in text or "\\sum" in text or "\\int" in text:
            return True
        # High density of math symbols
        math_chars = set("+=×÷∑∫∏√∞≤≥≠≈∂∇αβγδεζηθικλμνξπρστυφχψω")
        ratio = sum(1 for c in text if c in math_chars) / max(len(text), 1)
        return ratio > 0.2

    def _text_to_markdown_table(self, text: str) -> str:
        """Convert tab-separated text to markdown table."""
        lines = text.strip().split("\n")
        rows = [line.split("\t") for line in lines if line.strip()]
        return self._rows_to_markdown_table(rows)

    def _rows_to_markdown_table(self, rows: List[List[str]]) -> str:
        """Convert a list of rows to markdown table format."""
        if not rows:
            return ""

        # Ensure all rows have same number of columns
        max_cols = max(len(row) for row in rows)
        rows = [row + [""] * (max_cols - len(row)) for row in rows]

        # Build markdown
        header = "| " + " | ".join(rows[0]) + " |"
        separator = "| " + " | ".join(["---"] * max_cols) + " |"
        body_rows = []
        for row in rows[1:]:
            body_rows.append("| " + " | ".join(row) + " |")

        parts = [header, separator] + body_rows
        return "\n".join(parts)
